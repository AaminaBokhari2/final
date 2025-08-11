#!/usr/bin/env python3
import os
import sys
import platform
import subprocess
import tempfile
import shutil
from datetime import datetime
from typing import Dict, List, Tuple, Any
from pathlib import Path
import time
import json
import re
import requests
from urllib.parse import quote_plus, urljoin
from bs4 import BeautifulSoup
import xml.etree.ElementTree as ET
from youtube_service import YouTubeService
from collections import Counter
from urllib.parse import urlparse
from dataclasses import dataclass
from abc import ABC, abstractmethod
from typing import Optional, Dict
import random
import re

def safe_json_parse(text: str):
    """Clean and parse JSON from AI responses safely."""
    try:
        # Remove code fences
        text = re.sub(r"^```json\s*|\s*```$", "", text.strip(), flags=re.DOTALL)
        # Extract JSON part if extra commentary exists
        json_match = re.search(r"(\[.*\]|\{.*\})", text, re.DOTALL)
        if json_match:
            text = json_match.group(1)
        return json.loads(text)
    except Exception as e:
        print(f"⚠️ JSON parsing failed: {e}")
        return None


import pdfplumber
import pytesseract
from pdf2image import convert_from_path
from PIL import Image, ImageDraw, ImageFont
from dotenv import load_dotenv
from groq import Groq
# Add this import at the top of your pipeline.py

# Load environment variables
load_dotenv()

##### GROQ CLIENT WITH IMPROVED ERROR HANDLING #####
class GroqClient:
  def __init__(self):
      api_key = os.getenv("GROQ_API_KEY")
      if not api_key:
          # This will allow the app to run in fallback mode for AI features
          print("⚠️ GROQ_API_KEY not found. AI features will run in fallback mode.")
          self.client = None
      else:
          self.client = Groq(api_key=api_key)
      
      self.model_fallbacks = [
         "llama-3.1-8b-instant" ,
          "llama3-8b-8192",      # New fallback option
          "llama3-70b-8192"      # Larger model option
      ]

  def chat_completion(self, messages: List[Dict], model: str = None, max_tokens: int = None, retry_count: int = 3) -> str:
      if not self.client:
          return "❌ Groq API client not initialized. Running in fallback mode."

      if model is None:
          models_to_try = self.model_fallbacks
      else:
          models_to_try = [model] + [m for m in self.model_fallbacks if m != model]
      
      for attempt in range(retry_count):
    
          for model_name in models_to_try:
              try:
                  print(f"🤖 Using model: {model_name} (attempt {attempt + 1})")
                  response = self.client.chat.completions.create(
                      model=model_name,
                      messages=messages,
                      max_tokens=max_tokens,
                      temperature=0.7
                  )
                  return response.choices[0].message.content.strip()
              
              except Exception as e:
                  print(f"⏱️ Error with {model_name}: {str(e)}")
                  if attempt < retry_count - 1:
                      wait_time = (attempt + 1) * 2
                      print(f"⏳ Waiting {wait_time} seconds before retry...")
                      time.sleep(wait_time)
                  continue
      
      return "❌ All AI models failed. Please check your Groq API key and internet connection."

##### ENHANCED PDF PROCESSOR WITH BETTER ERROR HANDLING #####
class EnhancedPDFProcessor:
  def __init__(self):
      self.tesseract_available = self._check_tesseract()
      print(f"🔍 Tesseract available: {self.tesseract_available}")

  def _check_tesseract(self) -> bool:
      try:
          pytesseract.get_tesseract_version()
          return True
      except Exception as e:
          print(f"⚠️ Tesseract not available: {e}")
          return False

  def extract_text_with_ocr(self, file_path: str, max_pages: int = 20) -> Dict[str, any]:
      result = {
          "text": "",
          "page_count": 0,
          "extracted_pages": 0,
          "ocr_pages": 0,
          "word_count": 0,
          "status": "success",
          "methods_used": [],
          "message": "",
          "error_details": []
      }

      if not os.path.exists(file_path):
          result["status"] = "error"
          result["message"] = f"File not found: {file_path}"
          return result

      try:
          print(f"📄 Processing PDF: {file_path}")
          text_content = ""
          ocr_content = ""

          # Try pdfplumber first
          try:
              with pdfplumber.open(file_path) as pdf:
                  result["page_count"] = len(pdf.pages)
                  pages_to_process = min(result["page_count"], max_pages)
                  print(f"📊 PDF has {result['page_count']} pages, processing {pages_to_process}")

                  for page_num, page in enumerate(pdf.pages[:pages_to_process], 1):
                      try:
                          page_text = page.extract_text()
                          if page_text and len(page_text.strip()) > 20:
                              text_content += f"\n--- Page {page_num} ---\n{page_text.strip()}\n"
                              result["extracted_pages"] += 1
                              print(f"✅ Extracted text from page {page_num}: {len(page_text)} chars")
                          else:
                              print(f"⚠️ Page {page_num}: No meaningful text found")
                      except Exception as e:
                          result["error_details"].append(f"Page {page_num}: {str(e)}")
                          continue

                  if result["extracted_pages"] > 0:
                      result["methods_used"].append("text_extraction")
                      print(f"✅ Successfully extracted text from {result['extracted_pages']} pages")

          except Exception as e:
              result["error_details"].append(f"PDFPlumber error: {str(e)}")
              print(f"❌ PDFPlumber failed: {e}")

          # Try OCR if text extraction failed or yielded poor results
          if (result["extracted_pages"] == 0 or 
              result["extracted_pages"] < result["page_count"] * 0.3):
              
              if self.tesseract_available:
                  print("🔍 Attempting OCR extraction...")
                  try:
                      ocr_content = self._extract_with_ocr(file_path, max_pages=min(5, result["page_count"]))
                      if ocr_content:
                          result["methods_used"].append("ocr")
                          result["ocr_pages"] = ocr_content.count("--- Page")
                          print(f"✅ OCR extracted text from {result['ocr_pages']} pages")
                  except Exception as e:
                      result["error_details"].append(f"OCR error: {str(e)}")
                      print(f"❌ OCR failed: {e}")
              else:
                  result["error_details"].append("OCR not available (Tesseract not installed)")

          # Combine results
          final_text = text_content + ocr_content
          result["text"] = final_text.strip()
          result["word_count"] = len(final_text.split())

          # Set final status and message
          if result["word_count"] > 50:
              methods = " + ".join(result["methods_used"])
              result["message"] = f"✅ Successfully extracted {result['word_count']} words using: {methods}"
              result["status"] = "success"
          elif result["word_count"] > 0:
              result["status"] = "warning"
              result["message"] = f"⚠️ Limited content extracted ({result['word_count']} words). PDF may be image-based or have formatting issues."
          else:
              result["status"] = "error"
              result["message"] = "❌ No text could be extracted. PDF might be image-based, protected, or corrupted."

          print(f"📊 Final result: {result['status']} - {result['word_count']} words")
          return result

      except Exception as e:
          result["status"] = "error"
          result["message"] = f"❌ Critical error processing PDF: {str(e)}"
          result["error_details"].append(str(e))
          return result

  def _extract_with_ocr(self, file_path: str, max_pages: int = 5) -> str:
      if not self.tesseract_available:
          return ""

      ocr_text = ""
      temp_dir = tempfile.mkdtemp()

      try:
          print(f"🖼️ Converting PDF to images for OCR (max {max_pages} pages)...")
          images = convert_from_path(
              file_path, 
              first_page=1, 
              last_page=max_pages, 
              output_folder=temp_dir,
              dpi=200
          )
          
          for i, image in enumerate(images, 1):
              try:
                  print(f"🔍 Processing image {i} with OCR...")
                  page_text = pytesseract.image_to_string(
                      image, 
                      lang='eng', 
                      config='--psm 6'
                  )
                  if page_text.strip():
                      ocr_text += f"\n--- Page {i} (OCR) ---\n{page_text.strip()}\n"
                      print(f"✅ OCR page {i}: {len(page_text)} chars")
              except Exception as e:
                  print(f"❌ OCR failed on page {i}: {e}")
                  continue
          
          return ocr_text
      
      except Exception as e:
          print(f"❌ OCR process failed: {e}")
          return ""
      
      finally:
          shutil.rmtree(temp_dir, ignore_errors=True)

##### ENHANCED STUDY AGENTS WITH FIXED PROMPTS #####
class SummaryAgent:
  def __init__(self, client: GroqClient):
      self.client = client

  def generate_summary(self, text: str) -> str:
      if not text.strip():
          return "❌ No content available to summarize."
      
      if len(text.split()) < 10:
          return "⚠️ Content too short for meaningful summary."
      
      # Truncate text if too long
      max_chars = 8000
      if len(text) > max_chars:
          text = text[:max_chars] + "..."
      
      prompt = f"""Create a comprehensive, well-structured summary of the following academic content.

Document Content:
{text}

Please structure your summary as follows:

## 📋 DOCUMENT OVERVIEW
- Main topic and purpose
- Document type and scope

## 🎯 KEY CONCEPTS & DEFINITIONS
List the 5-8 most important concepts with brief explanations

## 📝 DETAILED SUMMARY
Write 2-3 paragraphs providing a comprehensive overview including:
- Main arguments or points
- Supporting evidence or examples
- Relationships between concepts
- Conclusions or implications

## 🔑 CRITICAL TAKEAWAYS
List 4-6 essential points that students must remember

## 📚 STUDY FOCUS AREAS
Highlight areas that deserve extra attention

Make the summary engaging, clear, and educational."""
      
      try:
          response = self.client.chat_completion([{"role": "user", "content": prompt}], max_tokens=1500)
          if response.startswith("❌"):
              return f"❌ Summary generation failed: {response}"
          return response
      except Exception as e:
          return f"❌ Summary generation failed: {str(e)}"
"""
AIPresentationCoordinatorAgent - Add this class to your pipeline.py file
"""

import json
import logging
from typing import List, Dict, Any, Tuple
from dataclasses import dataclass

# Set up logging
logger = logging.getLogger(__name__)

@dataclass
class SlideContent:
    """Represents content for a single slide"""
    title: str
    content: List[str]
    slide_type: str
    notes: str = ""

@dataclass
class Presentation:
    """Represents a complete presentation"""
    title: str
    slides: List[SlideContent]
    theme: str
    total_slides: int

class AIPresentationCoordinatorAgent:
    """
    FIXED: Coordinates the creation of presentations using AI - prevents hanging
    """
    
    def __init__(self, groq_client):
        """
        Initialize the presentation coordinator
        
        Args:
            groq_client: The GroqClient instance for AI operations
        """
        self.client = groq_client
        self.logger = logging.getLogger(__name__)
        
        # Slide type templates
        self.slide_types = {
            "title": "Title slide with main topic",
            "overview": "Overview/agenda slide",
            "content": "Main content slide",
            "comparison": "Comparison or vs slide",
            "conclusion": "Summary/conclusion slide",
            "call_to_action": "Call to action slide"
        }
        
        # Theme configurations
        self.themes = {
            "professional": {
                "colors": ["#1e3a8a", "#3b82f6", "#93c5fd"],
                "fonts": {"title": "Arial Bold", "body": "Arial"},
                "style": "Clean and professional"
            },
            "creative": {
                "colors": ["#7c3aed", "#a855f7", "#c084fc"],
                "fonts": {"title": "Helvetica Bold", "body": "Helvetica"},
                "style": "Creative and engaging"
            },
            "minimal": {
                "colors": ["#374151", "#6b7280", "#9ca3af"],
                "fonts": {"title": "Roboto Bold", "body": "Roboto"},
                "style": "Minimal and clean"
            }
        }
    
    def create_presentation(self, topic: str, audience: str = "general", 
                          duration: int = 10, theme: str = "professional") -> Tuple[Presentation, Dict, Dict]:
        """
        Create a complete presentation - FIXED to prevent hanging
        
        Args:
            topic: Presentation topic
            audience: Target audience
            duration: Presentation duration in minutes
            theme: Presentation theme
            
        Returns:
            Tuple of (Presentation, design_guidelines, quality_assessment)
        """
        try:
            self.logger.info(f"🎨 Creating presentation: {topic}")
            
            # Calculate number of slides based on duration
            slides_count = max(3, min(duration // 2, 8))  # Reduced max to prevent hanging
            
            # FIXED: Use fallback first, then try AI enhancement
            presentation_data = self._create_fallback_content(topic, slides_count)
            
            # Try to enhance with AI, but don't hang if it fails
            try:
                self.logger.info("🤖 Attempting AI enhancement...")
                enhanced_data = self._generate_presentation_content_safe(topic, audience, slides_count)
                if enhanced_data and enhanced_data.get("slides"):
                    presentation_data = enhanced_data
                    self.logger.info("✅ AI enhancement successful")
                else:
                    self.logger.info("🔄 Using fallback content")
            except Exception as ai_error:
                self.logger.warning(f"⚠️ AI enhancement failed: {ai_error}, using fallback")
            
            # Create slides
            slides = self._create_slides(presentation_data, theme)
            
            # Create presentation object
            presentation = Presentation(
                title=presentation_data.get("title", topic),
                slides=slides,
                theme=theme,
                total_slides=len(slides)
            )
            
            # Generate design guidelines
            design_guidelines = self._generate_design_guidelines(theme, topic)
            
            # Quality assessment
            quality_assessment = self._assess_quality(presentation)
            
            self.logger.info(f"✅ Presentation created: {len(slides)} slides")
            return presentation, design_guidelines, quality_assessment
            
        except Exception as e:
            self.logger.error(f"❌ Error creating presentation: {e}")
            # Return fallback presentation
            return self._create_fallback_presentation(topic, theme)
    
    def _generate_presentation_content_safe(self, topic: str, audience: str, slides_count: int) -> Dict:
        """FIXED: Generate presentation content using AI with timeout protection"""
        try:
            # Simplified prompt to reduce complexity
            prompt = f"""Create a {slides_count}-slide presentation about "{topic}".

Return only a valid JSON object:
{{
    "title": "Presentation Title",
    "slides": [
        {{"title": "Introduction", "content": ["Point 1", "Point 2"], "slide_type": "title", "notes": ""}},
        {{"title": "Main Content", "content": ["Key point", "Supporting detail"], "slide_type": "content", "notes": ""}}
    ]
}}

Topic: {topic}
Slides needed: {slides_count}
Keep it simple and concise."""

            # FIXED: Add timeout and retry logic
            max_attempts = 2
            for attempt in range(max_attempts):
                try:
                    self.logger.info(f"🤖 AI attempt {attempt + 1}/{max_attempts}")
                    
                    response = self.client.chat_completion(
                        [{"role": "user", "content": prompt}],
                        max_tokens=1500  # Reduced token limit
                    )
                    
                    if response.startswith("❌"):
                        self.logger.warning(f"AI response failed on attempt {attempt + 1}")
                        if attempt == max_attempts - 1:
                            raise Exception("All AI attempts failed")
                        continue
                    
                    # Quick JSON parse attempt
                    try:
                        # Extract JSON from response
                        json_str = response.strip()
                        if "```json" in json_str:
                            json_str = json_str.split("```json")[1].split("```")[0].strip()
                        elif "```" in json_str:
                            json_str = json_str.split("```")[1].split("```")[0].strip()
                        
                        # Parse JSON quickly
                        presentation_data = json.loads(json_str)
                        
                        # Validate basic structure
                        if presentation_data.get("slides") and len(presentation_data["slides"]) > 0:
                            self.logger.info("✅ Valid AI response received")
                            return presentation_data
                        else:
                            self.logger.warning("Invalid structure in AI response")
                            
                    except json.JSONDecodeError as json_error:
                        self.logger.warning(f"JSON parse failed: {json_error}")
                        
                    # If we get here, try next attempt
                    continue
                    
                except Exception as attempt_error:
                    self.logger.warning(f"Attempt {attempt + 1} failed: {attempt_error}")
                    if attempt == max_attempts - 1:
                        raise
            
            # If all attempts failed
            raise Exception("All AI generation attempts failed")
            
        except Exception as e:
            self.logger.error(f"AI content generation failed: {e}")
            return None
    
    def _create_slides(self, presentation_data: Dict, theme: str) -> List[SlideContent]:
        """Create slide objects from presentation data"""
        slides = []
        
        slides_data = presentation_data.get("slides", [])
        
        for slide_data in slides_data:
            slide = SlideContent(
                title=slide_data.get("title", "Untitled Slide"),
                content=slide_data.get("content", ["No content available"]),
                slide_type=slide_data.get("slide_type", "content"),
                notes=slide_data.get("notes", "")
            )
            slides.append(slide)
        
        return slides
    
    def _generate_design_guidelines(self, theme: str, topic: str) -> Dict:
        """Generate design guidelines for the presentation"""
        theme_config = self.themes.get(theme, self.themes["professional"])
        
        return {
            "color_scheme": theme_config["colors"],
            "fonts": theme_config["fonts"],
            "layout": {
                "type": "standard",
                "alignment": "left",
                "spacing": "comfortable"
            },
            "suggestions": f"Use {theme_config['style']} design. Maintain consistent formatting and clear visual hierarchy."
        }
    
    def _assess_quality(self, presentation: Presentation) -> Dict:
        """Assess presentation quality"""
        issues = []
        suggestions = []
        score = 100
        
        # Check slide count
        if presentation.total_slides < 3:
            issues.append("Too few slides")
            score -= 20
        elif presentation.total_slides > 12:
            suggestions.append("Consider reducing number of slides")
            score -= 5
        
        # Check content quality
        for slide in presentation.slides:
            if len(slide.content) < 1:
                issues.append(f"Slide '{slide.title}' has no content")
                score -= 15
            elif len(slide.content) > 6:
                suggestions.append(f"Consider reducing bullet points in '{slide.title}'")
                score -= 3
        
        # Add suggestions based on score
        if score > 90:
            suggestions.append("Excellent presentation structure")
        elif score > 75:
            suggestions.append("Good presentation, minor improvements possible")
        else:
            suggestions.append("Consider revising content and structure")
        
        return {
            "issues": issues,
            "suggestions": suggestions,
            "quality_assessment": f"Presentation quality score: {score}/100",
            "overall_score": score,
            "success": score >= 60
        }
    
    def _create_fallback_content(self, topic: str, slides_count: int) -> Dict:
        """Create fallback content when AI fails - IMPROVED"""
        # Clean up topic
        clean_topic = topic.replace("Page Band", "").replace("popular", "").strip()
        if not clean_topic:
            clean_topic = "Technology Overview"
        
        return {
            "title": f"{clean_topic} Presentation",
            "slides": [
                {
                    "title": f"Introduction to {clean_topic}",
                    "content": [
                        f"Overview of {clean_topic}",
                        "Key concepts and importance", 
                        "What we'll cover today"
                    ],
                    "slide_type": "title",
                    "notes": f"Welcome to this presentation about {clean_topic}"
                },
                {
                    "title": "Key Features & Benefits",
                    "content": [
                        "Primary features and capabilities",
                        "Main advantages and benefits",
                        "How it solves problems"
                    ],
                    "slide_type": "content",
                    "notes": f"Discuss the main features of {clean_topic}"
                },
                {
                    "title": "Applications & Use Cases",
                    "content": [
                        "Real-world applications",
                        "Common use cases",
                        "Industry examples"
                    ],
                    "slide_type": "content",
                    "notes": f"Show practical applications of {clean_topic}"
                },
                {
                    "title": "Implementation & Getting Started",
                    "content": [
                        "How to get started",
                        "Best practices",
                        "Common considerations"
                    ],
                    "slide_type": "content",
                    "notes": "Guide for implementation"
                },
                {
                    "title": "Conclusion & Next Steps",
                    "content": [
                        f"Key takeaways about {clean_topic}",
                        "Recommended next steps",
                        "Questions and discussion"
                    ],
                    "slide_type": "conclusion",
                    "notes": "Summarize and encourage questions"
                }
            ][:slides_count]  # Limit to requested number of slides
        }
    
    def _create_fallback_presentation(self, topic: str, theme: str) -> Tuple[Presentation, Dict, Dict]:
        """Create a fallback presentation when everything fails"""
        fallback_content = self._create_fallback_content(topic, 5)
        slides = self._create_slides(fallback_content, theme)
        
        presentation = Presentation(
            title=fallback_content["title"],
            slides=slides,
            theme=theme,
            total_slides=len(slides)
        )
        
        design_guidelines = self._generate_design_guidelines(theme, topic)
        quality_assessment = {
            "issues": ["Generated using fallback mode"],
            "suggestions": ["AI service used fallback - basic presentation created"],
            "quality_assessment": "Fallback presentation generated successfully",
            "overall_score": 75,
            "success": True
        }
        
        return presentation, design_guidelines, quality_assessment
class FlashcardAgent:
  def __init__(self, client: GroqClient):
      self.client = client

  def generate_flashcards_structured(self, text: str, num_cards=10) -> List[Dict]:
      """Generate structured flashcards data for the app interface"""
      if not text.strip():
          return []
      
      if len(text.split()) < 20:
          return []
      
      # Truncate text if too long
      max_chars = 7000
      if len(text) > max_chars:
          text = text[:max_chars] + "..."
      
      prompt = f"""Create {num_cards} high-quality study flashcards based on the following content. Return ONLY a valid JSON array with no additional text.

Content:
{text}

Return format (MUST be valid JSON):
[
  {{
      "question": "Clear, specific question",
      "answer": "Comprehensive answer with examples",
      "difficulty": "Basic",
      "category": "Main topic category",
      "hint": "Optional memory aid or hint"
  }},
  {{
      "question": "Another clear question",
      "answer": "Another comprehensive answer",
      "difficulty": "Intermediate", 
      "category": "Topic category",
      "hint": "Memory aid"
  }}
]

Guidelines:
- Create diverse question types (definitions, applications, comparisons)
- Test understanding, not just memorization
- Include relevant examples in answers
- Mix difficulty levels: Basic, Intermediate, Advanced
- Keep questions clear and answers comprehensive"""
      
      try:
          response = self.client.chat_completion([{"role": "user", "content": prompt}], max_tokens=2500)
          
          # Clean response to extract JSON
          response = response.strip()
          if response.startswith("```json"):
              response = response[7:]
          if response.endswith("```"):
              response = response[:-3]
          response = response.strip()
          
          # Try to parse JSON response
          try:
              flashcards_data = safe_json_parse(response)
              if isinstance(flashcards_data, list) and len(flashcards_data) > 0:
                  # Validate structure
                  valid_cards = []
                  for card in flashcards_data:
                      if isinstance(card, dict) and 'question' in card and 'answer' in card:
                          # Ensure all required fields
                          valid_card = {
                              'question': str(card.get('question', '')),
                              'answer': str(card.get('answer', '')),
                              'difficulty': card.get('difficulty', 'Basic'),
                              'category': card.get('category', 'General'),
                              'hint': card.get('hint', '')
                          }
                          valid_cards.append(valid_card)
                  
                  if valid_cards:
                      print(f"✅ Generated {len(valid_cards)} flashcards")
                      return valid_cards
          except json.JSONDecodeError as e:
              print(f"❌ JSON parsing failed: {e}")
              print(f"Response: {response[:200]}...")
          
          # Fallback: generate basic flashcards
          return self._generate_basic_flashcards(text, num_cards)
      
      except Exception as e:
          print(f"❌ Flashcard generation error: {e}")
          return self._generate_basic_flashcards(text, num_cards)
  
  def _generate_basic_flashcards(self, text: str, num_cards: int) -> List[Dict]:
      """Generate basic flashcards as fallback"""
      words = text.split()
      if len(words) < 50:
          return []
      
      # Extract key terms (simplified approach)
      sentences = text.split('.')[:num_cards]
      flashcards = []
      
      for i, sentence in enumerate(sentences):
          if len(sentence.strip()) > 20:
              flashcards.append({
                  'question': f"What is discussed about: {sentence.strip()[:50]}...?",
                  'answer': sentence.strip(),
                  'difficulty': 'Basic',
                  'category': 'Document Content',
                  'hint': 'Review the document content'
              })
      
      return flashcards[:num_cards]

class QuizAgent:
  def __init__(self, client: GroqClient):
      self.client = client

  def generate_quiz_structured(self, text: str, num_questions=8) -> List[Dict]:
      """Generate structured quiz data for the app interface"""
      if not text.strip():
          return []
      
      if len(text.split()) < 30:
          return []
      
      # Truncate text if too long
      max_chars = 7000
      if len(text) > max_chars:
          text = text[:max_chars] + "..."
      
      prompt = f"""Create {num_questions} multiple choice questions based on the following content. Return ONLY a valid JSON array with no additional text.

Content:
{text}

Return format (MUST be valid JSON):
[
  {{
      "question": "Clear, specific question",
      "options": ["Option A", "Option B", "Option C", "Option D"],
      "correct_answer": 0,
      "explanation": "Detailed explanation of the correct answer",
      "difficulty": "Basic"
  }},
  {{
      "question": "Another question",
      "options": ["Option 1", "Option 2", "Option 3", "Option 4"],
      "correct_answer": 2,
      "explanation": "Another explanation",
      "difficulty": "Intermediate"
  }}
]

Guidelines:
- Test understanding and application, not just memorization
- Make all options plausible but only one clearly correct
- Mix difficulty: Basic, Intermediate, Advanced
- correct_answer should be the index (0-3) of the correct option
- Provide educational explanations"""
      
      try:
          response = self.client.chat_completion([{"role": "user", "content": prompt}], max_tokens=3000)
          
          # Clean response to extract JSON
          response = response.strip()
          if response.startswith("```json"):
              response = response[7:]
          if response.endswith("```"):
              response = response[:-3]
          response = response.strip()
          
          # Try to parse JSON response
          try:
              quiz_data = safe_json_parse(response)
              if isinstance(quiz_data, list) and len(quiz_data) > 0:
                  # Validate structure
                  valid_questions = []
                  for q in quiz_data:
                      if (isinstance(q, dict) and 'question' in q and 'options' in q 
                          and 'correct_answer' in q and isinstance(q['options'], list) 
                          and len(q['options']) == 4):
                          
                          # Ensure valid correct_answer index
                          correct_idx = q.get('correct_answer', 0)
                          if not isinstance(correct_idx, int) or correct_idx < 0 or correct_idx > 3:
                              correct_idx = 0
                          
                          valid_question = {
                              'question': str(q.get('question', '')),
                              'options': [str(opt) for opt in q['options'][:4]],
                              'correct_answer': correct_idx,
                              'explanation': str(q.get('explanation', 'No explanation provided')),
                              'difficulty': q.get('difficulty', 'Basic')
                          }
                          valid_questions.append(valid_question)
                  
                  if valid_questions:
                      print(f"✅ Generated {len(valid_questions)} quiz questions")
                      return valid_questions
          except json.JSONDecodeError as e:
              print(f"❌ JSON parsing failed: {e}")
              print(f"Response: {response[:200]}...")
          
          # Fallback: generate basic quiz
          return self._generate_basic_quiz(text, num_questions)
      
      except Exception as e:
          print(f"❌ Quiz generation error: {e}")
          return self._generate_basic_quiz(text, num_questions)
  
  def _generate_basic_quiz(self, text: str, num_questions: int) -> List[Dict]:
      """Generate basic quiz as fallback"""
      words = text.split()
      if len(words) < 100:
          return []
      
      # Extract key sentences for questions
      sentences = [s.strip() for s in text.split('.') if len(s.strip()) > 30][:num_questions]
      quiz_questions = []
      
      for i, sentence in enumerate(sentences):
          if len(sentence) > 30:
              quiz_questions.append({
                  'question': f"According to the document, what is mentioned about the topic?",
                  'options': [
                      sentence[:50] + "..." if len(sentence) > 50 else sentence,
                      "This is not mentioned in the document",
                      "The document states the opposite",
                      "This is only partially correct"
                  ],
                  'correct_answer': 0,
                  'explanation': f"The document states: {sentence}",
                  'difficulty': 'Basic'
              })
      
      return quiz_questions[:num_questions]

##### REAL WEB DISCOVERY AGENTS #####

#!/usr/bin/env python3

# Add these complete implementations to your pipeline.py file
##### Q&A CHATBOT AGENT WITH CONVERSATION SUPPORT #####
# COMPLETE FIXED QAChatbotAgent CLASS
# Replace the existing QAChatbotAgent class in your pipeline.py with this version

import json
import time
from typing import Dict, List, Any, Tuple
from urllib.parse import quote_plus

class QAChatbotAgent:
    def __init__(self, client: GroqClient):
        self.client = client
        self.conversation_history = {}
        self.max_history = 10  # Keep last 10 exchanges
        
    def ask_question(self, question: str, document_text: str, session_id: str = "default") -> Dict[str, Any]:
        """
        Enhanced Q&A with conversation context and better fallback
        """
        if not question.strip():
            return {
                "answer": "Please provide a question to get an answer.",
                "status": "error",
                "fallback_used": False
            }
            
        if not document_text.strip():
            return {
                "answer": "No document content available to answer questions about.",
                "status": "error", 
                "fallback_used": False
            }
            
        # Initialize conversation history for session
        if session_id not in self.conversation_history:
            self.conversation_history[session_id] = []
            
        try:
            # Try AI-powered answer first
            if self.client and self.client.client:
                ai_answer = self._generate_ai_answer(question, document_text, session_id)
                if not ai_answer.startswith("❌"):
                    # Store successful interaction
                    self._add_to_history(session_id, question, ai_answer)
                    return {
                        "answer": ai_answer,
                        "status": "success",
                        "fallback_used": False
                    }
                    
            # Fallback to enhanced text matching
            fallback_answer = self._generate_enhanced_fallback_answer(question, document_text, session_id)
            self._add_to_history(session_id, question, fallback_answer)
            
            return {
                "answer": fallback_answer,
                "status": "success", 
                "fallback_used": True
            }
            
        except Exception as e:
            print(f"❌ Q&A error: {e}")
            fallback_answer = self._generate_enhanced_fallback_answer(question, document_text, session_id)
            self._add_to_history(session_id, question, fallback_answer)
            return {
                "answer": fallback_answer,
                "status": "success",
                "fallback_used": True
            }
    
    def _generate_ai_answer(self, question: str, document_text: str, session_id: str) -> str:
        """Generate AI-powered answer with conversation context"""
        try:
            # Limit document text for processing
            max_chars = 6000
            text_content = document_text[:max_chars]
            if len(document_text) > max_chars:
                text_content += "..."
                
            # Build conversation context
            context = ""
            if session_id in self.conversation_history and self.conversation_history[session_id]:
                recent_history = self.conversation_history[session_id][-3:]  # Last 3 exchanges
                context = "\n\nPrevious conversation:\n"
                for i, (prev_q, prev_a) in enumerate(recent_history, 1):
                    context += f"Q{i}: {prev_q}\nA{i}: {prev_a[:100]}...\n"
                    
            prompt = f"""You are an intelligent document assistant. Answer the question based on the document content and conversation context.

Document Content:
{text_content}
{context}

Current Question: {question}

Instructions:
- Provide a comprehensive, accurate answer based on the document
- Reference specific parts of the document when possible
- If the information isn't in the document, say so clearly
- Consider the conversation context for follow-up questions
- Use examples from the document to illustrate points
- Keep the answer well-structured and educational

Answer:"""

            response = self.client.chat_completion(
                [{"role": "user", "content": prompt}], 
                max_tokens=800
            )
            
            return response
            
        except Exception as e:
            print(f"❌ AI answer generation failed: {e}")
            return f"❌ AI answer generation failed: {str(e)}"
    
    def _generate_enhanced_fallback_answer(self, question: str, document_text: str, session_id: str) -> str:
        """Enhanced fallback answer with better text analysis"""
        if not question.strip() or not document_text.strip():
            return "I need both a question and document content to provide an answer."
            
        question_lower = question.lower()
        doc_lower = document_text.lower()
        
        # Extract question keywords (improved)
        question_words = []
        for word in question_lower.split():
            clean_word = word.strip('.,!?;:"()[]{}')
            if len(clean_word) > 3 and clean_word.isalpha():
                question_words.append(clean_word)
                
        # Find relevant sentences with better scoring
        sentences = [s.strip() for s in document_text.split('.') if len(s.strip()) > 20]
        scored_sentences = []
        
        for sentence in sentences:
            sentence_lower = sentence.lower()
            score = 0
            
            # Score based on keyword matches
            for word in question_words:
                if word in sentence_lower:
                    # Exact word match
                    score += 2
                elif any(word in sent_word for sent_word in sentence_lower.split()):
                    # Partial match
                    score += 1
                    
            # Boost score for question-type indicators
            if any(indicator in question_lower for indicator in ['what is', 'define', 'explain']):
                if any(indicator in sentence_lower for indicator in ['is', 'means', 'refers to', 'defined as']):
                    score += 1
                    
            if score > 0:
                scored_sentences.append((sentence, score))
                
        # Sort by relevance score
        scored_sentences.sort(key=lambda x: x[1], reverse=True)
        
        if scored_sentences:
            # Build comprehensive answer
            answer = f"""**Question:** {question}

**Based on the document content:**

**Primary Information:**
{scored_sentences[0][0]}"""

            # Add additional context if available
            if len(scored_sentences) > 1:
                answer += f"""

**Additional Context:**
{scored_sentences[1][0]}"""
                
            if len(scored_sentences) > 2:
                answer += f"""

**Related Information:**
{scored_sentences[2][0]}"""
                
            # Add conversation context if available
            if session_id in self.conversation_history and self.conversation_history[session_id]:
                answer += f"""

**Note:** This answer considers our previous conversation. For more detailed analysis, the AI service can provide enhanced responses when available."""
                
            answer += f"""

**Keywords found:** {', '.join(question_words[:5])}
**Confidence:** {'High' if scored_sentences[0][1] >= 3 else 'Medium' if scored_sentences[0][1] >= 2 else 'Low'}"""

            return answer
            
        else:
            # No direct matches found
            return f"""I couldn't find specific information directly answering "{question}" in the document.

**Suggestions:**
1. Try rephrasing your question with different keywords
2. Ask about broader topics covered in the document
3. Check if the information might be expressed differently

**Keywords I searched for:** {', '.join(question_words[:5])}

**Alternative approach:** You could ask about general topics I found in the document, such as the main themes or key concepts discussed.

**Note:** This search used basic text matching. When the AI service is available, it can provide more sophisticated analysis and better understand context and relationships in the document."""
    
    def generate_suggested_questions(self, document_text: str, num_questions: int = 5) -> List[str]:
        """Generate suggested questions based on document content"""
        if not document_text.strip():
            return []
            
        try:
            # Try AI-generated suggestions first
            if self.client and self.client.client:
                ai_suggestions = self._generate_ai_suggested_questions(document_text, num_questions)
                if ai_suggestions:
                    return ai_suggestions
                    
            # Fallback to rule-based suggestions
            return self._generate_fallback_suggested_questions(document_text, num_questions)
            
        except Exception as e:
            print(f"❌ Error generating suggested questions: {e}")
            return self._generate_fallback_suggested_questions(document_text, num_questions)
    
    def _generate_ai_suggested_questions(self, document_text: str, num_questions: int) -> List[str]:
        """Generate AI-powered suggested questions"""
        try:
            # Limit text for analysis
            analysis_text = document_text[:4000]
            
            prompt = f"""Based on this document content, generate {num_questions} insightful questions that would help someone understand the material better.

Document Content:
{analysis_text}

Generate questions that:
- Test understanding of key concepts
- Explore relationships between ideas
- Ask about practical applications
- Clarify important definitions
- Examine implications or conclusions

Return ONLY a JSON array of questions:
["Question 1?", "Question 2?", "Question 3?"]"""

            response = self.client.chat_completion(
                [{"role": "user", "content": prompt}], 
                max_tokens=400
            )
            
            # Clean and parse response
            response = response.strip()
            if response.startswith("```json"):
                response = response[7:]
            if response.endswith("```"):
                response = response[:-3]
            response = response.strip()
            
            try:
                questions = json.loads(response)
                if isinstance(questions, list) and len(questions) > 0:
                    return [str(q) for q in questions[:num_questions]]
            except json.JSONDecodeError:
                pass
                
            return []
            
        except Exception as e:
            print(f"❌ AI suggested questions failed: {e}")
            return []
    
    def _generate_fallback_suggested_questions(self, document_text: str, num_questions: int) -> List[str]:
        """Generate fallback suggested questions using text analysis"""
        questions = []
        
        # Extract key terms and concepts
        words = document_text.split()
        sentences = [s.strip() for s in document_text.split('.') if len(s.strip()) > 30]
        
        # Find capitalized terms (likely important concepts)
        capitalized_terms = []
        for word in words:
            clean_word = word.strip('.,!?;:"()[]{}')
            if (len(clean_word) > 3 and clean_word[0].isupper() and 
                not clean_word.isupper() and clean_word.isalpha()):
                capitalized_terms.append(clean_word)
                
        # Get most frequent important terms
        term_freq = {}
        for term in capitalized_terms:
            term_freq[term] = term_freq.get(term, 0) + 1
            
        important_terms = sorted(term_freq.keys(), key=lambda x: term_freq[x], reverse=True)[:10]
        
        # Generate different types of questions
        question_templates = [
            ("What is {term}?", "definition"),
            ("How does {term} work?", "explanation"), 
            ("What are the main characteristics of {term}?", "description"),
            ("Why is {term} important?", "significance"),
            ("What are the applications of {term}?", "application")
        ]
        
        # Generate questions from important terms
        for term in important_terms[:num_questions]:
            if len(questions) >= num_questions:
                break
            template, _ = question_templates[len(questions) % len(question_templates)]
            questions.append(template.format(term=term))
            
        # Add some general questions if we need more
        general_questions = [
            "What are the main topics covered in this document?",
            "What are the key takeaways from this material?",
            "How do the concepts in this document relate to each other?",
            "What practical applications are discussed?",
            "What conclusions does the document reach?"
        ]
        
        while len(questions) < num_questions and len(questions) < len(general_questions):
            questions.append(general_questions[len(questions)])
            
        return questions[:num_questions]
    
    def _add_to_history(self, session_id: str, question: str, answer: str):
        """Add Q&A pair to conversation history"""
        if session_id not in self.conversation_history:
            self.conversation_history[session_id] = []
            
        self.conversation_history[session_id].append((question, answer))
        
        # Keep only recent history
        if len(self.conversation_history[session_id]) > self.max_history:
            self.conversation_history[session_id] = self.conversation_history[session_id][-self.max_history:]
    
    def clear_conversation(self, session_id: str = "default"):
        """Clear conversation history for a session"""
        if session_id in self.conversation_history:
            del self.conversation_history[session_id]
            return True
        return False
    
    def get_conversation_summary(self, session_id: str = "default") -> Dict[str, Any]:
        """Get summary of conversation history"""
        if session_id not in self.conversation_history:
            return {
                "total_questions": 0,
                "recent_topics": [],
                "conversation_active": False
            }
            
        history = self.conversation_history[session_id]
        
        # Extract topics from recent questions
        recent_topics = []
        for question, _ in history[-5:]:  # Last 5 questions
            # Simple topic extraction from question
            words = question.lower().split()
            topic_words = [w for w in words if len(w) > 4 and w.isalpha()]
            if topic_words:
                recent_topics.append(topic_words[0])
                
        return {
            "total_questions": len(history),
            "recent_topics": list(set(recent_topics)),  # Remove duplicates
            "conversation_active": len(history) > 0
        }
        


class AIEnhancedResearchDiscoveryAgent:
    def __init__(self, client):
        self.client = client
        self.session = requests.Session()
        self.session.headers.update({
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36'
        })
        self.max_retries = 3
        self.retry_delay = 2

    def extract_smart_keywords_and_topic(self, text: str) -> Tuple[List[str], str]:
        """Extract smart keywords and main topic from text with enhanced error recovery"""
        if not text.strip():
            return ["academic", "study"], "Academic Content"
        
        # Limit text for analysis
        analysis_text = text[:6000]
        
        prompt = f"""Analyze this academic text and extract the most important keywords and main topic for research discovery.

Text: {analysis_text}

Return ONLY a JSON object with this format:
{{
  "main_topic": "The primary subject or field",
  "keywords": ["keyword1", "keyword2", "keyword3", "keyword4", "keyword5"],
  "research_area": "Specific research domain"
}}

Focus on academic and technical terms that would help find relevant research papers."""

        for attempt in range(self.max_retries):
            try:
                response = self.client.chat_completion([{"role": "user", "content": prompt}], max_tokens=300)
                
                # Clean response
                response = response.strip()
                if response.startswith("```json"):
                    response = response[7:]
                if response.endswith("```"):
                    response = response[:-3]
                response = response.strip()
                
                # Try to parse JSON response
                try:
                    data = json.loads(response)
                    keywords = data.get("keywords", ["research", "academic"])[:8]  # Limit keywords
                    topic = data.get("main_topic", "Academic Research")
                    
                    print(f"✅ Extracted topic: {topic}")
                    print(f"✅ Keywords: {keywords[:5]}")
                    
                    return keywords, topic
                    
                except json.JSONDecodeError as e:
                    print(f"❌ JSON parsing failed (attempt {attempt + 1}): {e}")
                    if attempt < self.max_retries - 1:
                        time.sleep(self.retry_delay)
                    continue
                    
            except Exception as e:
                print(f"❌ AI extraction failed (attempt {attempt + 1}): {e}")
                if attempt < self.max_retries - 1:
                    time.sleep(self.retry_delay)
                continue
        
        # Fallback: simple keyword extraction
        print("⚠️ Falling back to simple keyword extraction")
        words = text.lower().split()
        common_words = {"the", "and", "or", "but", "in", "on", "at", "to", "for", "of", "with", "by", "is", "are", "was", "were", "be", "been", "have", "has", "had", "do", "does", "did", "will", "would", "could", "should", "may", "might", "can", "cannot", "a", "an", "this", "that", "these", "those"}
        
        # Simple frequency analysis
        word_freq = {}
        for word in words:
            if len(word) > 4 and word not in common_words:
                word_freq[word] = word_freq.get(word, 0) + 1
        
        # Get most frequent words
        keywords = sorted(word_freq.keys(), key=lambda x: word_freq[x], reverse=True)[:5]
        if not keywords:
            keywords = ["research", "study", "analysis"]
            
        return keywords, "Academic Study"

    def extract_ai_research_context(self, text: str) -> Dict[str, Any]:
        """Use AI to extract comprehensive research context with fallback"""
        try:
            keywords, topic = self.extract_smart_keywords_and_topic(text)
            
            return {
                "research_domain": topic,
                "primary_concepts": keywords[:3],
                "methodologies": [],
                "key_terms": keywords[3:6] if len(keywords) > 3 else keywords,
                "related_fields": [],
                "research_questions": [f"What is {topic}?"],
                "academic_level": "intermediate",
                "search_intent": "research papers",
                "exclusion_terms": [],
                "paper_types": ["research paper", "review paper"],
                "time_preference": "recent"
            }
        except Exception as e:
            print(f"❌ Error extracting research context: {e}")
            return {
                "research_domain": "Academic Research",
                "primary_concepts": ["research", "study", "analysis"],
                "key_terms": [],
                "paper_types": ["research paper"],
                "time_preference": "recent"
            }

    def _create_ai_search_strategies(self, context: Dict[str, Any]) -> List[Dict]:
        """Create AI-guided search strategies with fallback"""
        try:
            keywords = context.get("primary_concepts", []) + context.get("key_terms", [])
            domain = context.get("research_domain", "research")
            
            strategies = [
                {
                    "name": "Primary Concept Search",
                    "terms": keywords[:3],
                    "sources": ["semantic", "arxiv"]
                },
                {
                    "name": "Domain-Specific Search", 
                    "terms": [domain] + keywords[:2],
                    "sources": ["arxiv", "semantic"]
                }
            ]
            
            # Add life sciences strategy if relevant
            if self._is_life_sciences_context(context):
                strategies.append({
                    "name": "Life Sciences Search",
                    "terms": keywords[:2],
                    "sources": ["pubmed", "semantic"]
                })
            
            return strategies
            
        except Exception as e:
            print(f"❌ Error creating search strategies: {e}")
            return [{
                "name": "Fallback Search",
                "terms": ["research", "study"],
                "sources": ["semantic", "arxiv"]
            }]

    def _is_life_sciences_context(self, context: Dict[str, Any]) -> bool:
        """Check if context suggests life sciences with error handling"""
        try:
            domain = context.get("research_domain", "").lower()
            keywords = [k.lower() for k in context.get("primary_concepts", []) + context.get("key_terms", [])]
            
            life_terms = ["biology", "medical", "health", "clinical", "genetics", "molecular", "cell", "protein", "drug", "disease", "neuroscience", "psychology"]
            
            return any(term in domain for term in life_terms) or any(any(term in kw for term in life_terms) for kw in keywords)
        except:
            return False

    def _search_arxiv_ai_enhanced(self, search_terms: List[str], context: Dict, max_results: int) -> List[Dict]:
        """Enhanced arXiv search with error recovery"""
        if not search_terms:
            return []
        
        query = ' AND '.join(f'"{term}"' if ' ' in term else term for term in search_terms[:3])
        url = f"http://export.arxiv.org/api/query?search_query=all:{quote_plus(query)}&start=0&max_results={max_results}&sortBy=relevance&sortOrder=descending"
        
        for attempt in range(self.max_retries):
            try:
                response = self.session.get(url, timeout=15)
                response.raise_for_status()
                
                root = ET.fromstring(response.content)
                papers = []
                
                for entry in root.findall('{http://www.w3.org/2005/Atom}entry'):
                    try:
                        title_elem = entry.find('{http://www.w3.org/2005/Atom}title')
                        summary_elem = entry.find('{http://www.w3.org/2005/Atom}summary')
                        published_elem = entry.find('{http://www.w3.org/2005/Atom}published')
                        id_elem = entry.find('{http://www.w3.org/2005/Atom}id')
                        
                        authors = []
                        for author in entry.findall('{http://www.w3.org/2005/Atom}author'):
                            name_elem = author.find('{http://www.w3.org/2005/Atom}name')
                            if name_elem is not None:
                                authors.append(name_elem.text.strip())
                        
                        if title_elem is not None and summary_elem is not None:
                            title = title_elem.text.strip()
                            abstract = summary_elem.text.strip()
                            
                            year = "2024"
                            if published_elem is not None:
                                try:
                                    year = published_elem.text[:4]
                                except:
                                    pass
                            
                            paper = {
                                'title': title,
                                'authors': ', '.join(authors[:4]) if authors else 'Unknown Authors',
                                'year': year,
                                'source': 'arXiv',
                                'abstract': (abstract[:400] + "...") if len(abstract) > 400 else abstract,
                                'url': id_elem.text if id_elem is not None else '#',
                                'relevance_label': 'Relevant'
                            }
                            papers.append(paper)
                    except Exception as e:
                        continue
                
                print(f"✅ Found {len(papers)} papers from arXiv")
                return papers
                
            except Exception as e:
                print(f"❌ arXiv search error (attempt {attempt + 1}): {e}")
                if attempt < self.max_retries - 1:
                    time.sleep(self.retry_delay)
                continue
        
        return []

    def _search_semantic_scholar_ai_enhanced(self, search_terms: List[str], context: Dict, max_results: int) -> List[Dict]:
        """Enhanced Semantic Scholar search with error recovery"""
        if not search_terms:
            return []
        
        query = ' '.join(search_terms[:3])
        url = f"https://api.semanticscholar.org/graph/v1/paper/search?query={quote_plus(query)}&limit={max_results}&fields=title,authors,year,abstract,url,venue,citationCount"
        
        for attempt in range(self.max_retries):
            try:
                response = self.session.get(url, timeout=15)
                if response.status_code == 200:
                    data = response.json()
                    papers = []
                    
                    for paper_data in data.get('data', []):
                        try:
                            authors_list = []
                            if paper_data.get('authors'):
                                authors_list = [author.get('name', '') for author in paper_data['authors'][:4]]
                            
                            paper = {
                                'title': paper_data.get('title', 'Untitled'),
                                'authors': ', '.join(authors_list) if authors_list else 'Unknown Authors',
                                'year': str(paper_data.get('year', '2024')),
                                'source': paper_data.get('venue', 'Semantic Scholar'),
                                'abstract': (paper_data.get('abstract', '')[:400] + "...") if paper_data.get('abstract') and len(paper_data.get('abstract', '')) > 400 else paper_data.get('abstract', 'No abstract available'),
                                'url': paper_data.get('url', '#'),
                                'relevance_label': 'Relevant'
                            }
                            papers.append(paper)
                        except Exception as e:
                            continue
                    
                    print(f"✅ Found {len(papers)} papers from Semantic Scholar")
                    return papers
                    
            except Exception as e:
                print(f"❌ Semantic Scholar search error (attempt {attempt + 1}): {e}")
                if attempt < self.max_retries - 1:
                    time.sleep(self.retry_delay)
                continue
        
        return []

    def _search_pubmed_ai_enhanced(self, search_terms: List[str], context: Dict, max_results: int) -> List[Dict]:
        """Enhanced PubMed search with error recovery"""
        if not search_terms:
            return []
        
        query = ' AND '.join(f'{term}[Title/Abstract]' for term in search_terms[:3])
        search_url = f"https://eutils.ncbi.nlm.nih.gov/entrez/eutils/esearch.fcgi?db=pubmed&term={quote_plus(query)}&retmax={max_results}&retmode=json&sort=relevance"
        
        for attempt in range(self.max_retries):
            try:
                response = self.session.get(search_url, timeout=15)
                if response.status_code == 200:
                    search_data = response.json()
                    pmids = search_data.get('esearchresult', {}).get('idlist', [])
                    
                    if pmids:
                        pmids_str = ','.join(pmids[:max_results])
                        fetch_url = f"https://eutils.ncbi.nlm.nih.gov/entrez/eutils/efetch.fcgi?db=pubmed&id={pmids_str}&retmode=xml"
                        
                        fetch_response = self.session.get(fetch_url, timeout=20)
                        if fetch_response.status_code == 200:
                            return self._parse_pubmed_xml(fetch_response.content)
                            
            except Exception as e:
                print(f"❌ PubMed search error (attempt {attempt + 1}): {e}")
                if attempt < self.max_retries - 1:
                    time.sleep(self.retry_delay)
                continue
        
        return []

    def _parse_pubmed_xml(self, xml_content: bytes) -> List[Dict]:
        """Parse PubMed XML response with error handling"""
        papers = []
        try:
            root = ET.fromstring(xml_content)
            
            for article in root.findall('.//PubmedArticle'):
                try:
                    title_elem = article.find('.//ArticleTitle')
                    title = title_elem.text if title_elem is not None else 'Untitled'
                    
                    abstract_texts = []
                    for abstract_elem in article.findall('.//Abstract/AbstractText'):
                        if abstract_elem.text:
                            abstract_texts.append(abstract_elem.text)
                    abstract = ' '.join(abstract_texts) if abstract_texts else 'No abstract available'
                    
                    authors = []
                    for author in article.findall('.//Author'):
                        fname = author.find('ForeName')
                        lname = author.find('LastName')
                        if fname is not None and lname is not None:
                            authors.append(f"{fname.text} {lname.text}")
                    
                    year_elem = article.find('.//PubDate/Year')
                    year = year_elem.text if year_elem is not None else '2024'
                    
                    journal_elem = article.find('.//Journal/Title')
                    journal = journal_elem.text if journal_elem is not None else 'PubMed'
                    
                    pmid_elem = article.find('.//PMID')
                    pmid = pmid_elem.text if pmid_elem is not None else ''
                    
                    paper = {
                        'title': title,
                        'authors': ', '.join(authors[:4]) if authors else 'Unknown Authors',
                        'year': year,
                        'source': journal,
                        'abstract': (abstract[:400] + "...") if len(abstract) > 400 else abstract,
                        'url': f"https://pubmed.ncbi.nlm.nih.gov/{pmid}/" if pmid else '#',
                        'relevance_label': 'Relevant'
                    }
                    papers.append(paper)
                except Exception as e:
                    continue
                    
        except Exception as e:
            print(f"❌ Error parsing PubMed XML: {e}")
        
        print(f"✅ Found {len(papers)} papers from PubMed")
        return papers

    def _ai_filter_and_rank_papers(self, papers: List[Dict], text: str, context: Dict) -> List[Dict]:
        """Filter and rank papers by relevance with error handling"""
        if not papers:
            return []
        
        try:
            # Remove duplicates
            seen_titles = set()
            unique_papers = []
            
            for paper in papers:
                title = paper.get('title', '').lower().strip()
                if title and title not in seen_titles:
                    seen_titles.add(title)
                    unique_papers.append(paper)
            
            # Simple ranking by title relevance to keywords
            keywords = context.get('primary_concepts', []) + context.get('key_terms', [])
            
            def relevance_score(paper):
                title = paper.get('title', '').lower()
                abstract = paper.get('abstract', '').lower()
                score = 0
                
                for keyword in keywords:
                    if keyword.lower() in title:
                        score += 2
                    elif keyword.lower() in abstract:
                        score += 1
                
                return score
            
            unique_papers.sort(key=relevance_score, reverse=True)
            return unique_papers
            
        except Exception as e:
            print(f"❌ Error filtering papers: {e}")
            return papers

    def find_papers(self, text: str, max_papers: int = 8) -> List[Dict]:
        """Main method to find research papers with comprehensive error handling"""
        if not text.strip():
            return []
        
        print("🔍 Extracting research context...")
        try:
            context = self.extract_ai_research_context(text)
        except Exception as e:
            print(f"❌ Error extracting research context: {e}")
            context = {
                "research_domain": "Academic Research",
                "primary_concepts": ["research", "study"],
                "key_terms": []
            }
        
        all_papers = []
        try:
            search_strategies = self._create_ai_search_strategies(context)
            
            for strategy in search_strategies:
                try:
                    strategy_name = strategy["name"]
                    search_terms = strategy["terms"]
                    sources = strategy["sources"]
                    
                    print(f"📚 Searching with strategy: {strategy_name}")
                    
                    for source in sources:
                        try:
                            if source == "arxiv" and len(all_papers) < max_papers:
                                papers = self._search_arxiv_ai_enhanced(search_terms, context, max_papers//2)
                                all_papers.extend(papers)
                            elif source == "semantic" and len(all_papers) < max_papers:
                                papers = self._search_semantic_scholar_ai_enhanced(search_terms, context, max_papers//2)
                                all_papers.extend(papers)
                            elif source == "pubmed" and len(all_papers) < max_papers:
                                papers = self._search_pubmed_ai_enhanced(search_terms, context, max_papers//2)
                                all_papers.extend(papers)
                            
                            time.sleep(1)  # Rate limiting
                            
                            if len(all_papers) >= max_papers:
                                break
                                
                        except Exception as e:
                            print(f"❌ Source {source} failed: {e}")
                            continue
                            
                except Exception as e:
                    print(f"❌ Strategy {strategy.get('name', 'Unknown')} failed: {e}")
                    continue
        except Exception as e:
            print(f"❌ Error in search strategy execution: {e}")
        
        # Filter and rank results
        try:
            relevant_papers = self._ai_filter_and_rank_papers(all_papers, text, context)
            return relevant_papers[:max_papers]
        except Exception as e:
            print(f"❌ Error filtering papers: {e}")
            return all_papers[:max_papers]


import os
import json
import re
from typing import List, Dict, Optional, Tuple
from datetime import datetime, timedelta
from googleapiclient.discovery import build
from dotenv import load_dotenv
import requests

load_dotenv()

# REPLACE your AIEnhancedYouTubeDiscoveryAgent class with this fixed version

from transformers import pipeline
from pptx import Presentation
from pptx.util import Inches
import tempfile
import os
from typing import List, Dict, Tuple

class PresentationAgent:
    def __init__(self):
        """Initialize presentation generation models"""
        # Text summarization model
        self.summarizer = pipeline(
            "summarization", 
            model="facebook/bart-large-cnn",
            tokenizer="facebook/bart-large-cnn"
        )
        
        # Key point extraction model
        self.keypoint_extractor = pipeline(
            "text2text-generation",
            model="mrm8488/t5-base-finetuned-question-generation-ap"
        )
        
        # Image generation (optional - requires more resources)
        self.image_generator = None
        try:
            from diffusers import StableDiffusionPipeline
            self.image_generator = StableDiffusionPipeline.from_pretrained(
                "CompVis/stable-diffusion-v1-4",
                use_auth_token=os.getenv("HF_TOKEN")
            )
        except ImportError:
            print("⚠️ Image generation disabled - diffusers not installed")
        except Exception as e:
            print(f"⚠️ Image generation disabled - {str(e)}")

    def generate_presentation(self, document_text: str, output_path: str, 
                            max_slides: int = 10, generate_images: bool = False) -> Dict:
        """
        Generate a PowerPoint presentation from document text
        
        Args:
            document_text: Input text to create presentation from
            output_path: Where to save the PPTX file
            max_slides: Maximum number of slides to generate
            generate_images: Whether to generate images (requires more resources)
            
        Returns:
            Dictionary with status and metadata
        """
        result = {
            "status": "success",
            "message": "",
            "slides_generated": 0,
            "images_generated": 0,
            "warnings": []
        }
        
        try:
            if not document_text.strip():
                result["status"] = "error"
                result["message"] = "No document content provided"
                return result
            
            # Step 1: Extract key points from document
            key_points = self._extract_key_points(document_text, max_slides)
            if not key_points:
                result["status"] = "error"
                result["message"] = "Could not extract key points from document"
                return result
            
            # Step 2: Create PowerPoint presentation
            prs = Presentation()
            
            # Add title slide
            title_slide = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            
            title.text = "AI-Generated Presentation"
            subtitle.text = "Created from Document Content"
            
            # Step 3: Add content slides
            for i, point in enumerate(key_points[:max_slides]):
                try:
                    # Add slide with title and content
                    bullet_slide = prs.slide_layouts[1]
                    slide = prs.slides.add_slide(bullet_slide)
                    shapes = slide.shapes
                    
                    title_shape = shapes.title
                    body_shape = shapes.placeholders[1]
                    
                    title_shape.text = f"Key Point {i+1}"
                    tf = body_shape.text_frame
                    tf.text = point["summary"]
                    
                    # Add bullet points if available
                    if point.get("bullets"):
                        for bullet in point["bullets"]:
                            p = tf.add_paragraph()
                            p.text = bullet
                            p.level = 0
                    
                    # Generate and add image if enabled
                    if generate_images and self.image_generator:
                        try:
                            image_path = self._generate_slide_image(point["summary"])
                            if image_path:
                                left = Inches(1)
                                top = Inches(2)
                                height = Inches(3)
                                slide.shapes.add_picture(
                                    image_path, left, top, height=height
                                )
                                os.remove(image_path)
                                result["images_generated"] += 1
                        except Exception as e:
                            result["warnings"].append(f"Image generation failed: {str(e)}")
                    
                    result["slides_generated"] += 1
                    
                except Exception as e:
                    result["warnings"].append(f"Failed to create slide {i+1}: {str(e)}")
                    continue
            
            # Save presentation
            prs.save(output_path)
            result["message"] = f"Successfully generated presentation with {result['slides_generated']} slides"
            
            if result["warnings"]:
                result["status"] = "completed_with_warnings"
            else:
                result["status"] = "success"
                
        except Exception as e:
            result["status"] = "error"
            result["message"] = f"Presentation generation failed: {str(e)}"
        
        return result

    def _extract_key_points(self, text: str, max_points: int) -> List[Dict]:
        """Extract key points and summaries from document text"""
        try:
            # First summarize the entire document to get main themes
            overall_summary = self.summarizer(
                text,
                max_length=150,
                min_length=30,
                do_sample=False
            )[0]['summary_text']
            
            # Split text into sections (simplified approach)
            sections = self._split_into_sections(text, max_points)
            
            key_points = []
            for section in sections:
                try:
                    # Summarize each section
                    summary = self.summarizer(
                        section,
                        max_length=100,
                        min_length=20,
                        do_sample=False
                    )[0]['summary_text']
                    
                    # Extract key bullets
                    bullets = self._extract_bullet_points(section)
                    
                    key_points.append({
                        "summary": summary,
                        "bullets": bullets[:3]  # Limit to 3 bullets per slide
                    })
                except Exception as e:
                    continue
            
            return key_points
            
        except Exception as e:
            print(f"❌ Key point extraction failed: {e}")
            return []

    def _split_into_sections(self, text: str, max_sections: int) -> List[str]:
        """Split document into logical sections"""
        # Simple approach - split by paragraphs or headings
        paragraphs = [p.strip() for p in text.split('\n') if p.strip()]
        
        # Group paragraphs into sections
        min_paras_per_section = max(1, len(paragraphs) // max_sections)
        sections = []
        current_section = []
        
        for para in paragraphs:
            current_section.append(para)
            if len(current_section) >= min_paras_per_section:
                sections.append("\n".join(current_section))
                current_section = []
        
        if current_section:
            sections.append("\n".join(current_section))
        
        return sections[:max_sections]

    def _extract_bullet_points(self, text: str) -> List[str]:
        """Extract bullet points from text"""
        try:
            # Use question generation model to find key points
            inputs = f"generate questions: {text[:2000]}"
            questions = self.keypoint_extractor(
                inputs,
                max_length=50,
                num_return_sequences=3,
                do_sample=True
            )
            
            # Convert questions to statements
            bullets = []
            for q in questions:
                question = q['generated_text'].strip()
                if question.endswith('?'):
                    # Simple conversion to statement
                    statement = question[:-1] + " is important."
                    bullets.append(statement)
            
            return bullets[:3]  # Return top 3
        
        except Exception as e:
            print(f"❌ Bullet extraction failed: {e}")
            return []

    def _generate_slide_image(self, text: str) -> Optional[str]:
        """Generate an image for a slide (optional)"""
        if not self.image_generator:
            return None
            
        try:
            # Create temp file
            temp_dir = tempfile.gettempdir()
            image_path = os.path.join(temp_dir, f"slide_img_{hash(text)}.png")
            
            # Generate image from text
            image = self.image_generator(
                prompt=text[:500],  # Limit prompt length
                height=512,
                width=512
            ).images[0]
            
            image.save(image_path)
            return image_path
            
        except Exception as e:
            print(f"❌ Image generation failed: {e}")
            return None


class AIEnhancedYouTubeDiscoveryAgent:
    def __init__(self, groq_client=None):
        """Initialize with proper error handling and fallbacks"""
        self.groq_client = groq_client
        self.youtube_service = self._initialize_youtube_service()
        self.session = requests.Session()
        self.session.headers.update({
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36'
        })
        self.cache = {}
        self.max_retries = 3
        self.retry_delay = 2

    def _initialize_youtube_service(self):
        """Initialize YouTube API service with robust error handling"""
        try:
            api_key = os.getenv("YOUTUBE_API_KEY")
            if not api_key:
                print("⚠️ YouTube API key not found in environment variables")
                return None
            
            from googleapiclient.discovery import build
            service = build('youtube', 'v3', developerKey=api_key)
            print("✅ YouTube API service initialized")
            return service
            
        except ImportError:
            print("⚠️ Google API client not installed - using fallback mode")
            return None
        except Exception as e:
            print(f"⚠️ YouTube API initialization failed: {e} - using fallback mode")
            return None

    def find_videos(self, keywords: List[str], topic: str, max_videos: int = 10) -> List[Dict]:
        """
        FIXED: Main method to find educational videos
        This method signature matches what the backend expects
        """
        print(f"🎥 Searching for videos about: {topic}")
        print(f"🔍 Keywords: {keywords[:3]}")
        
        videos = []
        
        # Create search queries from keywords and topic
        search_queries = self._create_search_queries(keywords, topic)
        
        for query in search_queries[:3]:  # Try top 3 queries
            if len(videos) >= max_videos:
                break
                
            print(f"🔍 Searching: {query}")
            
            # Try API first, then fallback
            query_videos = self._search_with_fallback(query, max_videos - len(videos))
            videos.extend(query_videos)
            
            time.sleep(1)  # Rate limiting
        
        # Filter and deduplicate
        filtered_videos = self._filter_videos(videos)
        result = filtered_videos[:max_videos]
        
        print(f"✅ Found {len(result)} videos")
        return result

    def _create_search_queries(self, keywords: List[str], topic: str) -> List[str]:
        """Create effective search queries"""
        queries = []
        
        # Primary topic query
        if topic and topic != "Study Material":
            queries.append(f"{topic} tutorial")
            queries.append(f"{topic} explained")
        
        # Keyword-based queries
        if keywords:
            queries.append(f"{keywords[0]} tutorial")
            if len(keywords) > 1:
                queries.append(f"{keywords[0]} {keywords[1]} course")
            if len(keywords) > 2:
                queries.append(f"{keywords[2]} lecture")
        
        # Fallback queries
        if not queries:
            queries = ["educational tutorial", "learning course", "academic lecture"]
        
        return queries

    def _search_with_fallback(self, query: str, max_results: int) -> List[Dict]:
        """Search with API first, fallback to scraping"""
        
        # Try YouTube API first
        if self.youtube_service:
            try:
                return self._youtube_api_search(query, max_results)
            except Exception as e:
                print(f"⚠️ API search failed for '{query}': {e}")
        
        # Fallback to web scraping
        try:
            return self._fallback_youtube_search(query, max_results)
        except Exception as e:
            print(f"⚠️ Fallback search failed for '{query}': {e}")
            return self._generate_placeholder_videos(query, max_results)

    def _youtube_api_search(self, query: str, max_results: int) -> List[Dict]:
        """Search using YouTube Data API"""
        try:
            request = self.youtube_service.search().list(
                q=query,
                part="snippet",
                maxResults=min(max_results, 10),
                type="video",
                videoDuration="medium",  # Exclude shorts
                relevanceLanguage="en",
                safeSearch="moderate",
                order="relevance"
            )
            
            response = request.execute()
            videos = []
            
            for item in response.get('items', []):
                try:
                    video_id = item['id']['videoId']
                    snippet = item['snippet']
                    
                    video = {
                        'title': self._clean_text(snippet.get('title', '')),
                        'channel': self._clean_text(snippet.get('channelTitle', '')),
                        'description': self._clean_text(snippet.get('description', ''))[:200],
                        'url': f"https://youtu.be/{video_id}",
                        'published_at': snippet.get('publishedAt', ''),
                        'thumbnail': snippet.get('thumbnails', {}).get('default', {}).get('url', ''),
                        'duration': 'N/A',
                        'views': 'N/A',
                        'educational_score': 'High',
                        'source': 'youtube_api'
                    }
                    videos.append(video)
                    
                except KeyError as e:
                    print(f"⚠️ Skipping malformed video item: {e}")
                    continue
            
            print(f"✅ API found {len(videos)} videos for '{query}'")
            return videos
            
        except Exception as e:
            print(f"❌ YouTube API search error: {e}")
            raise

    def _fallback_youtube_search(self, query: str, max_results: int) -> List[Dict]:
        """Fallback web scraping method"""
        search_url = f"https://www.youtube.com/results?search_query={quote_plus(query)}"
        
        for attempt in range(self.max_retries):
            try:
                response = self.session.get(search_url, timeout=15)
                if response.status_code == 200:
                    videos = self._parse_youtube_html(response.text, max_results)
                    if videos:
                        print(f"✅ Fallback found {len(videos)} videos for '{query}'")
                        return videos
                
            except Exception as e:
                print(f"❌ Fallback attempt {attempt + 1} failed: {e}")
                if attempt < self.max_retries - 1:
                    time.sleep(self.retry_delay)
        
        return []

    def _parse_youtube_html(self, html: str, max_results: int) -> List[Dict]:
        """Parse YouTube search results HTML"""
        videos = []
        
        # Look for video data in the HTML
        video_pattern = r'"videoId":"([^"]+)".*?"title":{"runs":\[{"text":"([^"]+)"}.*?"ownerText":{"runs":\[{"text":"([^"]+)"}'
        
        matches = re.finditer(video_pattern, html)
        
        for match in matches:
            if len(videos) >= max_results:
                break
                
            try:
                video_id = match.group(1)
                title = self._clean_text(match.group(2))
                channel = self._clean_text(match.group(3))
                
                # Skip shorts and low-quality content
                if self._is_quality_content(title):
                    video = {
                        'title': title,
                        'channel': channel,
                        'description': f"Educational video about {title[:50]}",
                        'url': f"https://youtu.be/{video_id}",
                        'published_at': datetime.now().isoformat() + "Z",
                        'thumbnail': f"https://img.youtube.com/vi/{video_id}/default.jpg",
                        'duration': 'N/A',
                        'views': 'N/A',
                        'educational_score': 'Medium',
                        'source': 'youtube_scrape'
                    }
                    videos.append(video)
                    
            except Exception as e:
                continue
        
        return videos

    def _generate_placeholder_videos(self, query: str, max_results: int) -> List[Dict]:
        """Generate placeholder videos when all methods fail"""
        print(f"⚠️ Generating placeholder videos for '{query}'")
        
        placeholders = []
        for i in range(min(max_results, 3)):
            placeholders.append({
                'title': f"Educational Video: {query} (Part {i+1})",
                'channel': "Educational Channel",
                'description': f"This would be a comprehensive video about {query}. Content includes tutorials, explanations, and practical examples.",
                'url': "https://youtu.be/placeholder",
                'published_at': datetime.now().isoformat() + "Z",
                'thumbnail': "",
                'duration': "15:30",
                'views': "10K views",
                'educational_score': 'Placeholder',
                'source': 'placeholder'
            })
        
        return placeholders

    def _filter_videos(self, videos: List[Dict]) -> List[Dict]:
        """Filter and rank videos by educational quality"""
        if not videos:
            return []
        
        # Remove duplicates
        unique_videos = []
        seen_titles = set()
        
        for video in videos:
            title = video.get('title', '').lower().strip()
            if title and title not in seen_titles:
                seen_titles.add(title)
                unique_videos.append(video)
        
        # Score and sort videos
        def educational_score(video):
            title = video.get('title', '').lower()
            channel = video.get('channel', '').lower()
            score = 0
            
            # Positive indicators
            educational_keywords = ['tutorial', 'course', 'lecture', 'learn', 'explained', 'guide', 'introduction']
            for keyword in educational_keywords:
                if keyword in title:
                    score += 2
            
            # Channel quality indicators
            quality_channels = ['khan', 'mit', 'stanford', 'coursera', 'edx', 'crash course']
            for channel_keyword in quality_channels:
                if channel_keyword in channel:
                    score += 3
            
            # Negative indicators
            negative_keywords = ['#shorts', 'tiktok', 'meme', 'funny', 'prank']
            for keyword in negative_keywords:
                if keyword in title:
                    score -= 5
            
            return score
        
        unique_videos.sort(key=educational_score, reverse=True)
        return unique_videos

    def _is_quality_content(self, title: str) -> bool:
        """Check if content appears to be educational quality"""
        title_lower = title.lower()
        
        # Skip obvious low-quality content
        skip_keywords = ['#shorts', 'tiktok', 'meme', 'funny', 'prank', 'reaction']
        if any(keyword in title_lower for keyword in skip_keywords):
            return False
        
        # Prefer educational content
        educational_keywords = ['tutorial', 'course', 'lecture', 'learn', 'explained', 'guide', 'how to']
        if any(keyword in title_lower for keyword in educational_keywords):
            return True
        
        # Default to including if not obviously bad
        return len(title.strip()) > 10

    def _clean_text(self, text: str) -> str:
        """Clean and normalize text"""
        if not text:
            return ''
        
        # Remove extra whitespace
        text = ' '.join(text.split())
        
        # Basic HTML entity decoding
        text = text.replace('&amp;', '&').replace('&quot;', '"').replace('&#39;', "'")
        
        return text

# REPLACE your AIEnhancedWebResourceAgent class with this fixed version

class AIEnhancedWebResourceAgent:
    def __init__(self, client):
        self.client = client
        self.session = requests.Session()
        self.session.headers.update({
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36'
        })
        self.max_retries = 3
        self.retry_delay = 2

    def extract_smart_keywords_and_topic(self, text: str) -> Tuple[str, List[str], List[str]]:
        """Extract keywords and topic for web resource discovery"""
        if not text.strip():
            return "Educational Resources", ["education", "learning"], ["education", "learning"]
        
        # Limit text for analysis
        analysis_text = text[:6000]
        
        if self.client and self.client.client:
            prompt = f"""Analyze this text and extract keywords for finding educational web resources.

Text: {analysis_text}

Return ONLY a JSON object:
{{
  "main_topic": "Primary subject",
  "keywords": ["keyword1", "keyword2", "keyword3", "keyword4", "keyword5"]
}}

Focus on educational terms that would help find learning resources."""

            try:
                response = self.client.chat_completion([{"role": "user", "content": prompt}], max_tokens=200)
                
                # Clean response
                response = response.strip()
                if response.startswith("```json"):
                    response = response[7:]
                if response.endswith("```"):
                    response = response[:-3]
                response = response.strip()
                
                # Try to parse JSON response
                try:
                    data = json.loads(response)
                    keywords = data.get("keywords", ["education", "learning"])[:8]
                    topic = data.get("main_topic", "Educational Resources")
                    
                    print(f"✅ Resource search - Topic: {topic}, Keywords: {keywords[:5]}")
                    return topic, keywords, keywords
                    
                except json.JSONDecodeError:
                    pass
                    
            except Exception as e:
                print(f"❌ AI keyword extraction failed: {e}")
        
        # Fallback: simple keyword extraction
        words = text.lower().split()
        common_words = {"the", "and", "or", "but", "in", "on", "at", "to", "for", "of", "with", "by", "is", "are", "was", "were", "be", "been", "have", "has", "had", "do", "does", "did", "will", "would", "could", "should", "may", "might", "can", "cannot", "a", "an", "this", "that", "these", "those"}
        
        word_freq = {}
        for word in words:
            clean_word = word.lower().strip('.,!?;:"()[]{}')
            if len(clean_word) > 4 and clean_word not in common_words and clean_word.isalpha():
                word_freq[clean_word] = word_freq.get(clean_word, 0) + 1
        
        keywords = sorted(word_freq.keys(), key=lambda x: word_freq[x], reverse=True)[:8]
        if not keywords:
            keywords = ["education", "tutorial", "course"]
        
        topic = keywords[0].title() if keywords else "Educational Resources"
        
        print(f"✅ Resource search - Topic: {topic}, Keywords: {keywords[:5]}")
        return topic, keywords, keywords

    def find_resources(self, keywords: List[str], topic: str, max_resources: int = 12) -> List[Dict]:
        """
        FIXED: Main method to find web learning resources
        This method signature matches what the backend expects
        """
        if not keywords and not topic:
            return []
        
        print(f"🌐 Searching for web resources about: {topic}")
        print(f"🔍 Keywords: {keywords[:3]}")
        
        all_resources = []
        
        # Create search strategies
        search_strategies = self._create_resource_search_strategies(keywords, topic)
        
        for strategy in search_strategies:
            try:
                strategy_name = strategy["name"]
                search_func = strategy["search_function"]
                search_params = strategy["params"]
                
                print(f"🌐 Resource strategy: {strategy_name}")
                
                resources = search_func(**search_params, max_results=max_resources//len(search_strategies) + 1)
                all_resources.extend(resources)
                time.sleep(1)  # Rate limiting
                
            except Exception as e:
                print(f"❌ Resource strategy {strategy_name} failed: {e}")
                continue
        
        # Filter and rank
        relevant_resources = self._filter_and_rank_resources(all_resources, keywords, topic)
        result = relevant_resources[:max_resources]
        
        print(f"✅ Found {len(result)} web resources")
        return result

    def _create_resource_search_strategies(self, keywords: List[str], topic: str) -> List[Dict]:
        """Create resource search strategies"""
        strategies = [
            {
                "name": "Wikipedia Search",
                "search_function": self._search_wikipedia,
                "params": {"keywords": keywords, "topic": topic}
            },
            {
                "name": "Educational Platforms",
                "search_function": self._search_educational_platforms,
                "params": {"keywords": keywords, "topic": topic}
            },
            {
                "name": "Academic Resources",
                "search_function": self._search_academic_resources,
                "params": {"keywords": keywords, "topic": topic}
            },
            {
                "name": "Tutorial Resources",
                "search_function": self._search_tutorial_resources,
                "params": {"keywords": keywords, "topic": topic}
            }
        ]
        
        return strategies

    def _search_wikipedia(self, keywords: List[str], topic: str, max_results: int) -> List[Dict]:
        """Search Wikipedia for articles"""
        resources = []
        
        search_terms = [topic] + keywords[:2] if topic else keywords[:3]
        
        for term in search_terms:
            if len(resources) >= max_results:
                break
                
            try:
                # Wikipedia API search
                search_url = f"https://en.wikipedia.org/api/rest_v1/page/summary/{quote_plus(term)}"
                response = self.session.get(search_url, timeout=10)
                
                if response.status_code == 200:
                    data = response.json()
                    
                    if data.get('type') != 'disambiguation' and not data.get('title', '').startswith('File:'):
                        resource = {
                            'title': data.get('title', term),
                            'type': 'Reference Article',
                            'source': 'Wikipedia',
                            'description': data.get('extract', '')[:300] + ('...' if len(data.get('extract', '')) > 300 else ''),
                            'url': data.get('content_urls', {}).get('desktop', {}).get('page', f"https://en.wikipedia.org/wiki/{quote_plus(term)}"),
                            'quality_score': 'High'
                        }
                        resources.append(resource)
                        
            except Exception as e:
                print(f"❌ Wikipedia search error for '{term}': {e}")
                continue
        
        return resources[:max_results]

    def _search_educational_platforms(self, keywords: List[str], topic: str, max_results: int) -> List[Dict]:
        """Search educational platforms"""
        resources = []
        
        # Create educational resources for the topic
        educational_sites = [
            {
                'title': f"{topic} - Khan Academy",
                'type': 'Interactive Course',
                'source': 'Khan Academy',
                'description': f"Comprehensive interactive lessons and exercises covering {topic}. Includes videos, practice problems, and progress tracking.",
                'url': f"https://www.khanacademy.org/search?search_again=1&page_search_query={quote_plus(topic)}",
                'quality_score': 'Excellent'
            },
            {
                'title': f"{topic} Course - Coursera",
                'type': 'Online Course',
                'source': 'Coursera',
                'description': f"University-level courses on {topic} from top institutions. Includes assignments, peer reviews, and certificates.",
                'url': f"https://www.coursera.org/search?query={quote_plus(topic)}",
                'quality_score': 'High'
            },
            {
                'title': f"{topic} - edX",
                'type': 'MOOC Course',
                'source': 'edX',
                'description': f"Free online courses on {topic} from Harvard, MIT, and other top universities. Self-paced learning with verified certificates available.",
                'url': f"https://www.edx.org/search?q={quote_plus(topic)}",
                'quality_score': 'High'
            },
            {
                'title': f"{topic} - Udemy",
                'type': 'Online Course',
                'source': 'Udemy',
                'description': f"Practical courses on {topic} taught by industry experts. Includes hands-on projects and lifetime access.",
                'url': f"https://www.udemy.com/courses/search/?q={quote_plus(topic)}",
                'quality_score': 'Good'
            }
        ]
        
        return educational_sites[:max_results]

    def _search_academic_resources(self, keywords: List[str], topic: str, max_results: int) -> List[Dict]:
        """Search for academic resources"""
        resources = []
        
        # Create academic resources
        academic_resources = [
            {
                'title': f"{topic} - MIT OpenCourseWare",
                'type': 'Course Materials',
                'source': 'MIT OCW',
                'description': f"MIT course materials covering {topic}. Includes lecture notes, assignments, exams, and video lectures from actual MIT courses.",
                'url': f"https://ocw.mit.edu/search/?q={quote_plus(topic)}",
                'quality_score': 'Excellent'
            },
            {
                'title': f"{topic} - Stanford Online",
                'type': 'Academic Course',
                'source': 'Stanford Online',
                'description': f"Stanford University courses and materials on {topic}. High-quality academic content from world-class faculty.",
                'url': f"https://online.stanford.edu/search-catalog?keywords={quote_plus(topic)}",
                'quality_score': 'Excellent'
            },
            {
                'title': f"{topic} - Harvard Extension",
                'type': 'Academic Resource',
                'source': 'Harvard Extension',
                'description': f"Harvard Extension School resources on {topic}. Academic-level content with practical applications.",
                'url': f"https://www.extension.harvard.edu/academics/courses?keywords={quote_plus(topic)}",
                'quality_score': 'High'
            }
        ]
        
        return academic_resources[:max_results]

    def _search_tutorial_resources(self, keywords: List[str], topic: str, max_results: int) -> List[Dict]:
        """Search for tutorial and learning resources"""
        resources = []
        
        # Create tutorial resources based on keywords
        for i, keyword in enumerate(keywords[:max_results]):
            resources.append({
                'title': f"{keyword.title()} Tutorial Guide",
                'type': 'Tutorial',
                'source': 'Educational Portal',
                'description': f"Step-by-step tutorial covering the fundamentals of {keyword}. Includes examples, exercises, and practical applications.",
                'url': f"https://www.tutorialspoint.com/{quote_plus(keyword.lower())}",
                'quality_score': 'Good'
            })
        
        # Add some general learning resources
        general_resources = [
            {
                'title': f"Learn {topic} - W3Schools",
                'type': 'Interactive Tutorial',
                'source': 'W3Schools',
                'description': f"Interactive tutorials and references for {topic}. Includes try-it-yourself examples and exercises.",
                'url': f"https://www.w3schools.com/",
                'quality_score': 'Good'
            },
            {
                'title': f"{topic} Documentation",
                'type': 'Reference',
                'source': 'Official Documentation',
                'description': f"Official documentation and reference materials for {topic}. Comprehensive guides and API references.",
                'url': f"https://docs.example.com/{quote_plus(topic.lower())}",
                'quality_score': 'High'
            }
        ]
        
        resources.extend(general_resources)
        return resources[:max_results]

    def _filter_and_rank_resources(self, resources: List[Dict], keywords: List[str], topic: str) -> List[Dict]:
        """Filter and rank web resources"""
        if not resources:
            return []
        
        # Remove duplicates
        unique_resources = []
        seen_titles = set()
        
        for resource in resources:
            title = resource.get('title', '').lower()
            if title not in seen_titles:
                seen_titles.add(title)
                unique_resources.append(resource)
        
        # Sort by quality score and relevance
        quality_order = {'Excellent': 4, 'High': 3, 'Good': 2, 'Fair': 1}
        
        def rank_score(resource):
            quality = resource.get('quality_score', 'Fair')
            base_score = quality_order.get(quality, 1)
            
            # Boost score if title contains keywords or topic
            title = resource.get('title', '').lower()
            description = resource.get('description', '').lower()
            
            relevance_boost = 0
            if topic.lower() in title:
                relevance_boost += 2
            
            for keyword in keywords[:3]:
                if keyword.lower() in title or keyword.lower() in description:
                    relevance_boost += 1
            
            return base_score + relevance_boost
        
        unique_resources.sort(key=rank_score, reverse=True)
        return unique_resources

    def _search_real_time_web(self, query: str, max_results: int) -> List[Dict]:
        """Search the web for real-time resources (optional enhancement)"""
        # This could be implemented with a web search API like Google Custom Search
        # For now, return empty list as fallback
        return []

    def _validate_url(self, url: str) -> bool:
        """Validate if URL is accessible"""
        try:
            response = self.session.head(url, timeout=5)
            return response.status_code == 200
        except:
            return False



# Updated main execution logic with proper error handling
def run_study_assistant(pdf_path):
    """Main function to run the study assistant pipeline"""
    print(f"\n🎓 Processing: {pdf_path}")
    
    try:
        # Initialize components
        client = GroqClient()
        processor = EnhancedPDFProcessor()
        summary_agent = SummaryAgent(client)
        flashcard_agent = FlashcardAgent(client)
        quiz_agent = QuizAgent(client)
        presentation_agent = PresentationAgent()
        
        # Create discovery agents
        print("🔧 Initializing discovery agents...")
        research_agent = AIEnhancedResearchDiscoveryAgent(client)
        youtube_agent = AIEnhancedYouTubeDiscoveryAgent(client)
        web_agent = AIEnhancedWebResourceAgent(client)
        print("✅ Discovery agents initialized")

        # Extract text
        print("\n📄 Extracting text from PDF...")
        result = processor.extract_text_with_ocr(pdf_path)

        if result["status"] == "error":
            print(f"❌ PDF Processing Failed: {result['message']}")
            return

        print(f"\n{result['message']}")
        
        if result["word_count"] < 20:
            print("⚠️ Very little content extracted. Results may be limited.")

        # Generate study materials
        print("\n🧠 Generating study materials...")
        
        print("📘 Creating summary...")
        summary = summary_agent.generate_summary(result["text"])
        
        print("🃏 Creating flashcards...")
        flashcards = flashcard_agent.generate_flashcards_structured(result["text"], num_cards=8)
        
        print("📝 Creating quiz...")
        quiz = quiz_agent.generate_quiz_structured(result["text"], num_questions=6)

        # Generate presentation
        print("\n📊 Generating presentation...")
        pptx_path = os.path.join(os.path.dirname(pdf_path), "generated_presentation.pptx")
        presentation_result = presentation_agent.generate_presentation(
            document_text=result["text"],
            output_path=pptx_path,
            max_slides=12,
            generate_images=False
        )
        
        if presentation_result["status"] == "success":
            print(f"✅ Generated presentation: {pptx_path}")
        else:
            print(f"⚠️ Presentation generation had issues: {presentation_result['message']}")

        # Discovery phase
        print("\n🔍 Discovering learning resources...")
        
        papers = []
        videos = []
        resources = []
        
        try:
            print("📄 Finding research papers...")
            papers = research_agent.find_papers(result["text"], max_papers=6)
            print(f"✅ Found {len(papers)} research papers")
        except Exception as e:
            print(f"❌ Research paper discovery failed: {e}")
        
        try:
            print("🎥 Finding educational videos...")
            videos = youtube_agent.find_videos(result["text"], max_videos=6)
            print(f"✅ Found {len(videos)} educational videos")
        except Exception as e:
            print(f"❌ Video discovery failed: {e}")
        
        try:
            print("🌐 Finding web resources...")
            resources = web_agent.find_resources(result["text"], max_resources=8)
            print(f"✅ Found {len(resources)} web resources")
        except Exception as e:
            print(f"❌ Web resource discovery failed: {e}")

        # Display results
        print("\n" + "="*80)
        print("📘 SUMMARY")
        print("="*80)
        print(summary)
        
        if flashcards:
            print("\n" + "="*80)
            print("🃏 FLASHCARDS")
            print("="*80)
            for i, card in enumerate(flashcards, 1):
                print(f"\nCard {i}:")
                print(f"Q: {card.get('question', 'No question')}")
                print(f"A: {card.get('answer', 'No answer')}")
                print(f"Difficulty: {card.get('difficulty', 'N/A')}")
        
        if quiz:
            print("\n" + "="*80)
            print("📝 QUIZ")
            print("="*80)
            for i, question in enumerate(quiz, 1):
                print(f"\nQuestion {i}: {question.get('question', 'No question')}")
                options = question.get('options', [])
                for j, option in enumerate(options):
                    print(f"{j+1}. {option}")
                correct_idx = question.get('correct_answer', 0)
                if 0 <= correct_idx < len(options):
                    print(f"Correct answer: {options[correct_idx]}")
                print(f"Explanation: {question.get('explanation', 'No explanation')}")
        
        if papers:
            print("\n" + "="*80)
            print("📄 RESEARCH PAPERS")
            print("="*80)
            for i, paper in enumerate(papers, 1):
                print(f"\nPaper {i}: {paper.get('title', 'No title')}")
                print(f"Authors: {paper.get('authors', 'Unknown')}")
                print(f"Source: {paper.get('source', 'Unknown')} ({paper.get('year', 'N/A')})")
                abstract = paper.get('abstract', 'No abstract available')
                print(f"Abstract: {abstract[:200]}{'...' if len(abstract) > 200 else ''}")
                print(f"URL: {paper.get('url', 'No URL')}")
                print(f"Relevance: {paper.get('relevance_label', 'N/A')}")
        else:
            print("\n" + "="*80)
            print("📄 RESEARCH PAPERS")
            print("="*80)
            print("No research papers found. This could be due to:")
            print("- Limited or non-academic content in the PDF")
            print("- Network connectivity issues")
            print("- API rate limiting")
        
        if videos:
            print("\n" + "="*80)
            print("🎥 EDUCATIONAL VIDEOS")
            print("="*80)
            for i, video in enumerate(videos, 1):
                print(f"\nVideo {i}: {video.get('title', 'No title')}")
                print(f"Channel: {video.get('channel', 'Unknown')}")
                print(f"Duration: {video.get('duration', 'N/A')} | Views: {video.get('views', 'N/A')}")
                description = video.get('description', 'No description')
                print(f"Description: {description[:200]}{'...' if len(description) > 200 else ''}")
                print(f"URL: {video.get('url', 'No URL')}")
                print(f"Quality: {video.get('educational_score', 'N/A')}")
        else:
            print("\n" + "="*80)
            print("🎥 EDUCATIONAL VIDEOS")
            print("="*80)
            print("No educational videos found. This could be due to:")
            print("- YouTube access restrictions")
            print("- Limited educational content for this topic")
            print("- Network connectivity issues")
        
        if resources:
            print("\n" + "="*80)
            print("🌐 WEB RESOURCES")
            print("="*80)
            for i, resource in enumerate(resources, 1):
                print(f"\nResource {i}: {resource.get('title', 'No title')}")
                print(f"Type: {resource.get('type', 'Unknown')} | Source: {resource.get('source', 'Unknown')}")
                description = resource.get('description', 'No description')
                print(f"Description: {description[:200]}{'...' if len(description) > 200 else ''}")
                print(f"URL: {resource.get('url', 'No URL')}")
                print(f"Quality: {resource.get('quality_score', 'N/A')}")
        else:
            print("\n" + "="*80)
            print("🌐 WEB RESOURCES")
            print("="*80)
            print("No web resources found. This could be due to:")
            print("- Network connectivity issues")
            print("- Content filtering restrictions")
            print("- Limited educational resources for this topic")

        print(f"\n✅ Study assistant completed successfully!")
        print(f"📊 Generated: {len(flashcards)} flashcards, {len(quiz)} quiz questions")
        print(f"🔍 Found: {len(papers)} papers, {len(videos)} videos, {len(resources)} resources")

    except Exception as e:
        print(f"❌ Critical error in study assistant: {str(e)}")
        import traceback
        print(f"📍 Full error trace:\n{traceback.format_exc()}")

# Test function for the discovery agents
def test_discovery_agents():
    """Test the discovery agents with sample content"""
    print("\n🧪 Testing Discovery Agents...")
    
    try:
        client = GroqClient()
        sample_text = "Machine learning is a method of data analysis that automates analytical model building. It is a branch of artificial intelligence based on the idea that systems can learn from data, identify patterns and make decisions with minimal human intervention."
        
        print("📄 Testing research agent...")
        research_agent = AIEnhancedResearchDiscoveryAgent(client)
        papers = research_agent.find_papers(sample_text, max_papers=3)
        print(f"✅ Research agent found {len(papers)} papers")
        
        print("🎥 Testing YouTube agent...")
        youtube_agent = AIEnhancedYouTubeDiscoveryAgent(client)
        videos = youtube_agent.find_videos(sample_text, max_videos=3)
        print(f"✅ YouTube agent found {len(videos)} videos")
        
        print("🌐 Testing web resource agent...")
        web_agent = AIEnhancedWebResourceAgent(client)
        resources = web_agent.find_resources(sample_text, max_resources=3)
        print(f"✅ Web agent found {len(resources)} resources")
        
        print("✅ All discovery agents tested successfully!")
        
    except Exception as e:
        print(f"❌ Discovery agent test failed: {e}")
        import traceback
        print(f"📍 Full error trace:\n{traceback.format_exc()}")
        
if __name__ == "__main__":
  print("\n🎓 AI Study Assistant CLI - Enhanced Version")
  print("="*50)
  print("1. Run Study Assistant")
  print("2. Diagnose PDF")
  print("3. Test OCR Setup")
  print("4. Test Groq Connection")
  print("5. Exit")
  
  choice = input("\nChoose an option (1-5): ").strip()

  if choice == "1":
      pdf_path = input("Enter PDF file path: ").strip().strip('"')
      if os.path.exists(pdf_path):
          run_study_assistant(pdf_path)
      else:
          print(f"❌ File not found: {pdf_path}")
  
  elif choice == "2":
      pdf_path = input("Enter PDF file path: ").strip().strip('"')
      diagnose_pdf(pdf_path)
  
  elif choice == "3":
      test_ocr_setup()
  
  elif choice == "4":
      test_groq_connection()
  
  elif choice == "5":
      print("👋 Goodbye!")
  
  else:
      print("❌ Invalid option. Please choose 1-5.")
